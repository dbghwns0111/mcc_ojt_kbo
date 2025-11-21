from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 프레젠테이션 객체 생성
prs = Presentation()

# 슬라이드 레이아웃 정의 (보통 0:Title, 1:Title and Content, ...)
TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
BULLET_SLIDE_LAYOUT = prs.slide_layouts[1]

def add_slide(prs, layout, title_text, content_text_list=None, subtitle_text=None):
    """슬라이드를 추가하는 헬퍼 함수"""
    slide = prs.slides.add_slide(layout)
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = title_text

    # 서브타이틀 처리 (타이틀 슬라이드용)
    if subtitle_text and layout == TITLE_SLIDE_LAYOUT:
        subtitle = shapes.placeholders[1]
        subtitle.text = subtitle_text
    
    # 본문 불릿 포인트 처리
    if content_text_list and layout == BULLET_SLIDE_LAYOUT:
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for i, text in enumerate(content_text_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.level = 0 # 1차 불릿 레벨

            # 하위 불릿 포인트 처리 (들여쓰기) -> 여기서는 단순화를 위해 모두 1레벨로 처리하거나 수동 줄바꿈 이용

# --- 슬라이드 내용 작성 시작 ---

# Slide 1: 타이틀
add_slide(
    prs,
    TITLE_SLIDE_LAYOUT,
    title_text="KBO 데이터 분석 프로젝트 주간 보고",
    subtitle_text="데이터 파이프라인 구축 및 스키마 설계 현황\n\n발표자: 서비스 개발 부서 인턴 OOO\n날짜: 2025.MM.DD"
)

# Slide 2: 프로젝트 개요 및 목표
content_s2 = [
    "프로젝트 주제: 2021~2025년 KBO 데이터를 활용한 MSTR 기반 BI 대시보드 구축",
    "프로젝트 목표",
    "- (기술) 데이터 수집부터 DW, MSTR 모델링까지 E2E BI 프로세스 경험",
    "- (비즈니스) 리그 흥행 분석, 선수 가치 평가 등 경영진 관점 인사이트 도출",
    "현재 진행 단계: 데이터 수집 완료 → 전처리 완료 → [DB 스키마 구축 및 MSTR 모델링 중]"
]
add_slide(prs, BULLET_SLIDE_LAYOUT, "프로젝트 개요 및 목표", content_s2)

# Slide 3: 주요 성과 1 - 데이터 수집 및 현황
content_s3 = [
    "수집 기간: 최근 5개년 (2021년 ~ 2025년 현재)",
    "데이터 원천: KBO 공식 홈페이지 및 주요 기록 사이트 (크롤링)",
    "확보 데이터 범위(Scope)",
    "- 마스터: 선수 기본 정보, 구단 정보, 경기장 정보",
    "- 트랜잭션: 경기별 메타데이터 (날짜, 구장, 관중, 시간 등)",
    "- 성적: 타자/투수 일별 기록, 상황별 기록, 심화 지표(WAR 등)",
    "성과: 분석에 필요한 거시적/미시적 데이터를 모두 확보하여 기반 마련"
]
add_slide(prs, BULLET_SLIDE_LAYOUT, "주요 성과 1: 데이터 수집 현황", content_s3)

# Slide 4: 주요 성과 2 - 데이터 전처리 (ETL 핵심)
content_s4 = [
    "주요 이슈: 일별 선수 기록에 '경기ID' 부재로 메타데이터 연결 불가",
    "해결 방안 (핵심 성과)",
    "1. '연도별 선수 소속팀 이력' 데이터 추가 수집",
    "2. 파이썬 활용: [날짜 + 소속팀 + 상대팀] 조합으로 유니크 경기ID 매핑 성공",
    "결과: 모든 일별 기록이 특정 경기와 연결되어 다차원 분석 가능해짐",
    "기타: 날짜 형식 표준화, 결측치 처리(취소 경기), 데이터 타입 통일 완료"
]
add_slide(prs, BULLET_SLIDE_LAYOUT, "주요 성과 2: 데이터 전처리 (ETL)", content_s4)

# Slide 5: 주요 성과 3 - 스키마 구성안 (Star Schema)
content_s5 = [
    "모델링 전략: MSTR 성능 최적화 및 직관적 분석을 위한 '스타 스키마' 채택",
    "[스키마 구조 - ERD Placeholder]",
    "Dimension (분석 기준)",
    "- DIM_PLAYER (선수), DIM_GAME (경기/시계열), DIM_TEAM (구단 메타)",
    "Fact (측정값)",
    "- FACT_DAILY (일별 성적), FACT_ADVANCED (심화 지표)",
    "현황: MySQL DB 적재 완료, MSTR 스키마 오브젝트 생성 진행 중"
]
add_slide(prs, BULLET_SLIDE_LAYOUT, "주요 성과 3: 스키마 구성안", content_s5)
# *참고: 여기에 실제 ERD 이미지를 넣으려면 코드가 더 복잡해져서, 생성 후 PPT에서 직접 이미지를 넣는 것을 추천합니다.

# Slide 6: 향후 계획 - 대시보드 구상안
content_s6 = [
    "컨셉: One-Click Filtering 기반의 통합 인사이트 대시보드",
    "구성 계획 (4 Tabs)",
    "1. KBO Overview: 리그/구단 KPI, 성적, 흥행 추이 통합 (메인)",
    "2. Game Changer: 피치클락, ABS 등 신규 규정 도입 효과 분석",
    "3. Player Scout: '연봉 대비 WAR(가성비)' 분석을 통한 숨은 인재 발굴",
    "4. Prediction (GG): 데이터 기반 골든글러브 예측 및 1:1 비교 시뮬레이션"
]
add_slide(prs, BULLET_SLIDE_LAYOUT, "향후 계획: 대시보드 구상안", content_s6)

# Slide 7: Next Steps
content_s7 = [
    "차주 계획",
    "- MSTR 스키마 작업 완료 (관계 설정, 계층 구조 정의)",
    "- 기본 및 복합 메트릭(Metric) 생성 (타율, 방어율 등 계산식)",
    "- 'KBO Overview' 탭 시각화 프로토타입 구현 시작"
]
add_slide(prs, BULLET_SLIDE_LAYOUT, "Next Steps", content_s7)


# PPT 파일 저장
file_name = "KBO_주간보고.pptx"
prs.save(file_name)
print(f"'{file_name}' 파일이 성공적으로 생성되었습니다!")